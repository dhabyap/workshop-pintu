#!/usr/bin/env python3
"""Generate workshop PPTX — FROM ZERO TO TRADING BOTS."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_BG = RGBColor(15, 23, 42)
NAVY = RGBColor(26, 58, 92)
ACCENT = RGBColor(52, 211, 153)  # mint green
GOLD = RGBColor(251, 191, 36)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(150, 150, 160)
LIGHT_BG = RGBColor(30, 41, 59)

def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        from lxml import etree
    return shape

def add_text(slide, text, left, top, width, height, font_size=24, bold=False, color=WHITE, align=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tf

def add_para(tf, text, font_size=20, bold=False, color=GRAY, space_before=Pt(8), font_name='Calibri'):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.space_before = space_before
    return p

def add_code_block(slide, code_text, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(20, 30, 50)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_top = Pt(8)
    p = tf.paragraphs[0]
    p.text = code_text
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(52, 211, 153)
    p.font.name = 'Consolas'
    return tf

def slide_number(slide, num):
    add_text(slide, f"{num:02d}", Inches(12.3), Inches(6.9), Inches(1), Inches(0.5), font_size=12, color=GRAY, align=PP_ALIGN.RIGHT)

def section_divider(slide, top, color=ACCENT):
    add_shape(slide, Inches(0.8), top, Inches(1.5), Inches(0.06), color)

# ─── SLIDE 1: TITLE ───
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)
# Accent bar
add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), ACCENT)
# Big title
add_text(slide, "FROM ZERO", Inches(1.2), Inches(1.5), Inches(10), Inches(1.2), font_size=60, bold=True, color=WHITE)
add_text(slide, "TO TRADING BOTS", Inches(1.2), Inches(2.7), Inches(10), Inches(1.2), font_size=60, bold=True, color=ACCENT)
section_divider(slide, Inches(4.2))
add_text(slide, "Workshop PINTU x Web3 Dev Bandung", Inches(1.2), Inches(4.6), Inches(8), Inches(0.5), font_size=22, color=WHITE)
add_text(slide, "25 Juli 2026 | 13:00 WIB | Parla Cafe", Inches(1.2), Inches(5.2), Inches(8), Inches(0.5), font_size=16, color=GRAY)
add_text(slide, "Instructor: Putra", Inches(1.2), Inches(5.8), Inches(4), Inches(0.5), font_size=14, color=GOLD)
slide_number(slide, 1)

# ─── SLIDE 2: ABOUT ME ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, "About Me", Inches(0.8), Inches(0.5), Inches(6), Inches(0.8), font_size=36, bold=True, color=WHITE)
section_divider(slide, Inches(1.3))
items = [
    "Web developer sejak 2019",
    "Laravel, PHP, Web3",
    "Crypto + jewelry trader",
    "Jalanin XAU signal bot 24/7",
    "Bot real dengan indicator: RSI, EMA, MACD, BB, ATR",
]
y = Inches(2.0)
for item in items:
    add_text(slide, f"  >  {item}", Inches(1), y, Inches(11), Inches(0.5), font_size=20, color=GRAY)
    y += Inches(0.65)
add_text(slide, "Bot gue hari ini:", Inches(1), Inches(5.0), Inches(10), Inches(0.5), font_size=16, color=ACCENT)
add_text(slide, "XAU LONG active | Entry $4,019 | P&L +1.41% | WR 25% (1W/4L)", Inches(1), Inches(5.5), Inches(10), Inches(0.5), font_size=14, color=GOLD)
slide_number(slide, 2)

# ─── SLIDE 3: AGENDA ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, "Agenda (45 Menit)", Inches(0.8), Inches(0.5), Inches(8), Inches(0.8), font_size=36, bold=True, color=WHITE)
section_divider(slide, Inches(1.3))
agenda = [
    ("5m", "Install AI Assistant"),
    ("5m", "Setup Provider (OpenRouter gratis)"),
    ("5m", "Download Bot Template"),
    ("10m", "Jalanin Bot + Lihat Sinyal"),
    ("10m", "Cron — Bot Auto 24/7"),
    ("5m", "Bonus: Telegram Notifikasi"),
    ("5m", "Q&A"),
]
y = Inches(2.0)
for dur, title in agenda:
    add_text(slide, f"{dur}", Inches(1), y, Inches(1.2), Inches(0.5), font_size=18, bold=True, color=ACCENT)
    add_text(slide, title, Inches(2.5), y, Inches(8), Inches(0.5), font_size=18, color=WHITE)
    y += Inches(0.6)
slide_number(slide, 3)

# ─── SLIDE 4: KENAPA BOT SINYAL ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, "Kenapa Bot Sinyal?", Inches(0.8), Inches(0.5), Inches(8), Inches(0.8), font_size=36, bold=True, color=WHITE)
section_divider(slide, Inches(1.3))
add_text(slide, "Lo ga perlu jadi programmer — cukup:", Inches(1), Inches(1.8), Inches(10), Inches(0.5), font_size=20, color=GRAY)
items = ["Bisa buka terminal", "Bisa copy-paste", "Mau nyobain"]
y = Inches(2.5)
for item in items:
    add_text(slide, f"   [X] {item}", Inches(1.5), y, Inches(8), Inches(0.5), font_size=18, color=WHITE)
    y += Inches(0.6)
add_text(slide, "AI (Hermes Agent) nanganin sisanya:", Inches(1), Inches(4.5), Inches(10), Inches(0.5), font_size=18, color=ACCENT)
add_text(slide, "Download template -> Jalanin -> Cron -> Sinyal masuk HP", Inches(1.5), Inches(5.1), Inches(10), Inches(0.5), font_size=16, color=GOLD)
slide_number(slide, 4)

# ─── SLIDE 5: KENAPA HERMES ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, "Kenapa Hermes Agent?", Inches(0.8), Inches(0.5), Inches(10), Inches(0.8), font_size=36, bold=True, color=WHITE)
section_divider(slide, Inches(1.3))
features = [
    ("Open Source", "Gratis, ga ada biaya langganan"),
    ("Cron Bawaan", "Bot jalan otomatis tanpa server"),
    ("Multi-Provider", "Claude, GPT, DeepSeek, pilih bebas"),
    ("Multi-Platform", "Telegram, Discord, Email, File"),
]
y = Inches(2.0)
for title, desc in features:
    add_text(slide, title, Inches(1.2), y, Inches(4), Inches(0.5), font_size=22, bold=True, color=ACCENT)
    add_text(slide, desc, Inches(5.5), y, Inches(6), Inches(0.5), font_size=18, color=GRAY)
    y += Inches(1.0)
slide_number(slide, 5)

# ─── SLIDE 6: INSTALL ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, "Install Hermes Agent", Inches(0.8), Inches(0.5), Inches(10), Inches(0.8), font_size=36, bold=True, color=WHITE)
section_divider(slide, Inches(1.3))
add_text(slide, "Buka terminal, jalanin:", Inches(1), Inches(1.8), Inches(10), Inches(0.5), font_size=18, color=GRAY)
add_code_block(slide, "$ curl -fsSL https://hermes-agent.nousresearch.com/install.sh | bash", Inches(1), Inches(2.5), Inches(11), Inches(0.9))
add_text(slide, "Cek hasil:", Inches(1), Inches(3.8), Inches(10), Inches(0.5), font_size=18, color=GRAY)
add_code_block(slide, "$ hermes doctor", Inches(1), Inches(4.3), Inches(11), Inches(0.6))
add_text(slide, "[Windows] Python 3.10+ harus udah terinstall.", Inches(1.5), Inches(5.3), Inches(10), Inches(0.5), font_size=16, color=GOLD)
add_text(slide, "Bantu temen sebelah kalau error!", Inches(1.5), Inches(5.8), Inches(10), Inches(0.5), font_size=16, color=GOLD)
slide_number(slide, 6)

# ─── SLIDE 7: SETUP PROVIDER ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, "Setup Provider (Gratis)", Inches(0.8), Inches(0.5), Inches(10), Inches(0.8), font_size=36, bold=True, color=WHITE)
section_divider(slide, Inches(1.3))
add_text(slide, "Step 1:  Daftar OpenRouter (dari rumah)", Inches(1), Inches(1.8), Inches(11), Inches(0.5), font_size=18, color=GRAY)
add_text(slide, "Buka openrouter.ai/keys -> Login -> Create Key -> Copy API key", Inches(1.5), Inches(2.4), Inches(10), Inches(0.5), font_size=16, color=WHITE)
add_text(slide, "Step 2:  Setup di Hermes", Inches(1), Inches(3.2), Inches(11), Inches(0.5), font_size=18, color=GRAY)
add_code_block(slide, "$ hermes model", Inches(1), Inches(3.8), Inches(11), Inches(0.6))
add_text(slide, "Pilih: OpenRouter -> paste API key -> pilih DeepSeek V3", Inches(1.5), Inches(4.6), Inches(10), Inches(0.5), font_size=16, color=WHITE)
add_text(slide, "Rekomendasi: DeepSeek V3 (gratis, cepet)", Inches(1), Inches(5.3), Inches(10), Inches(0.5), font_size=16, bold=True, color=ACCENT)
add_text(slide, "Alternatif: Gemini API (60 req/menit) | Ollama local (butuh RAM 8GB+)", Inches(1), Inches(5.9), Inches(11), Inches(0.5), font_size=14, color=GRAY)
slide_number(slide, 7)

# ─── SLIDE 8: TEMPLATE BOT ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, "Template Bot Udah Siap", Inches(0.8), Inches(0.5), Inches(10), Inches(0.8), font_size=36, bold=True, color=WHITE)
section_divider(slide, Inches(1.3))
add_text(slide, "trading_bot_template.py  ~  230 baris", Inches(1), Inches(1.8), Inches(8), Inches(0.5), font_size=20, bold=True, color=ACCENT)
add_text(slide, "Isinya:", Inches(1), Inches(2.5), Inches(8), Inches(0.5), font_size=16, color=GRAY)
items = [
    "Fetch harga (demo mode, tinggal ganti API)",
    "5 Indikator: RSI, EMA, MACD, BB, ATR",
    "Generate sinyal LONG / SHORT + SL / TP",
    "Simpan history ke JSON",
    "Cooldown + silent mode (ga spam)",
]
y = Inches(3.0)
for item in items:
    add_text(slide, f"   [>] {item}", Inches(1.5), y, Inches(10), Inches(0.4), font_size=15, color=WHITE)
    y += Inches(0.5)
add_text(slide, "Ga perlu nulis kode. Tinggal download dan jalanin.", Inches(1), Inches(5.8), Inches(10), Inches(0.5), font_size=18, bold=True, color=GOLD)
slide_number(slide, 8)

# ─── SLIDE 9: JALANIN BOT ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, "Jalanin Bot", Inches(0.8), Inches(0.5), Inches(8), Inches(0.8), font_size=36, bold=True, color=WHITE)
section_divider(slide, Inches(1.3))
add_code_block(slide, "$ python trading_bot_template.py", Inches(1), Inches(1.8), Inches(11), Inches(0.6))
add_text(slide, "Output:", Inches(1), Inches(2.7), Inches(10), Inches(0.5), font_size=16, color=GRAY)
output = (
    "TRADING SIGNAL - 25 Jul 13:30 WIB\n"
    "========================================\n"
    "Price: $4,075.70 (+1.71%)\n"
    "Direction: LONG | Confidence: 95%\n"
    "Entry: $4,075.70 | SL: $3,962.50 | TP: $4,113.20\n"
    "R:R 1:1.67"
)
add_code_block(slide, output, Inches(1), Inches(3.0), Inches(11), Inches(2.0))
add_text(slide, "Bot juga simpan history di signals.json + history.json", Inches(1), Inches(5.4), Inches(10), Inches(0.5), font_size=14, color=GRAY)
slide_number(slide, 9)

# ─── SLIDE 10: BOT SILENT ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, "Kenapa Bot Diam?", Inches(0.8), Inches(0.5), Inches(10), Inches(0.8), font_size=36, bold=True, color=WHITE)
section_divider(slide, Inches(1.3))
add_text(slide, "Bot lo ga ngomong tiap detik. Itu FITUR, bukan bug.", Inches(1), Inches(1.8), Inches(11), Inches(0.5), font_size=20, bold=True, color=GOLD)
items = [
    ("Cooldown 4 jam", "Jeda antar sinyal biar ga spam"),
    ("Active guard", "Ga generate sinyal kalo masih ada posisi aktif"),
    ("Silent mode", "Output cuma kalo ada sinyal BARU / posisi CLOSED"),
]
y = Inches(2.6)
for title, desc in items:
    add_text(slide, title, Inches(1.2), y, Inches(4), Inches(0.5), font_size=18, bold=True, color=ACCENT)
    add_text(slide, desc, Inches(5.5), y, Inches(6), Inches(0.5), font_size=16, color=GRAY)
    y += Inches(0.9)
add_text(slide, "Mau liat status? Buka signals.json atau jalanin ulang", Inches(1), Inches(5.5), Inches(10), Inches(0.5), font_size=16, color=GRAY)
slide_number(slide, 10)

# ─── SLIDE 11: LIVE DEMO ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, "LIVE DEMO: Bot XAU Saya", Inches(0.8), Inches(0.5), Inches(10), Inches(0.8), font_size=36, bold=True, color=WHITE)
section_divider(slide, Inches(1.3))
demo = (
    "XAU SIGNAL - 21 Jul 23:25 WIB\n"
    "Active: LONG | Entry $4,019\n"
    "Current: $4,075.70 | P&L +1.41%\n"
    "SL: $3,962.50 | TP: $4,113.20\n"
    "[ ====---------------- ] 36/48h"
)
add_code_block(slide, demo, Inches(1), Inches(2.0), Inches(8), Inches(2.2))
add_text(slide, "WR: 25% (1W/4L) — tapi profit karena risk management.", Inches(1.2), Inches(4.5), Inches(10), Inches(0.5), font_size=18, color=GOLD)
add_text(slide, "Ini bot gue. Jalan 24/7, kirim sinyal ke Telegram.", Inches(1.2), Inches(5.1), Inches(10), Inches(0.5), font_size=18, color=WHITE)
add_text(slide, "Gue ga pernah coding ulang. Template -> Cron -> Done.", Inches(1.2), Inches(5.7), Inches(10), Inches(0.5), font_size=16, color=GRAY)
slide_number(slide, 11)

# ─── SLIDE 12: CRON ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, "Cron: Bot Jalan Otomatis", Inches(0.8), Inches(0.5), Inches(10), Inches(0.8), font_size=36, bold=True, color=WHITE)
section_divider(slide, Inches(1.3))
add_text(slide, "Manual python trading_bot.py tiap hari? Capek.", Inches(1), Inches(1.8), Inches(11), Inches(0.5), font_size=18, color=GRAY)
add_text(slide, "Cron = bot jalan sendiri, lo tinggal terima hasil.", Inches(1), Inches(2.4), Inches(11), Inches(0.5), font_size=20, bold=True, color=ACCENT)
add_code_block(slide, (
    "$ hermes cron create 'every 4h' \\\n"
    '    --name "Sinyal XAU Auto" \\\n'
    '    --prompt "Jalankan python trading_bot_template.py 2>&1. Kirim output apa adanya."'
), Inches(1), Inches(3.2), Inches(11), Inches(1.8))
add_text(slide, "Cek cron:", Inches(1), Inches(5.3), Inches(10), Inches(0.5), font_size=16, color=GRAY)
add_code_block(slide, "$ hermes cron list", Inches(1), Inches(5.7), Inches(5), Inches(0.5))
slide_number(slide, 12)

# ─── SLIDE 13: TELEGRAM ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, "Bonus: Telegram Notifikasi", Inches(0.8), Inches(0.5), Inches(10), Inches(0.8), font_size=36, bold=True, color=WHITE)
section_divider(slide, Inches(1.3))
add_text(slide, "Setup Gateway:", Inches(1), Inches(1.8), Inches(8), Inches(0.5), font_size=18, color=GRAY)
add_code_block(slide, "$ hermes gateway setup", Inches(1), Inches(2.3), Inches(8), Inches(0.6))
steps = [
    "1. Pilih Telegram",
    "2. Bikin bot lewat @BotFather",
    "3. Start bot di HP",
    "4. Update cron -> deliver ke Telegram",
]
y = Inches(3.3)
for s in steps:
    add_text(slide, s, Inches(1.5), y, Inches(8), Inches(0.4), font_size=16, color=WHITE)
    y += Inches(0.5)
add_code_block(slide, '$ hermes cron update <job_id> --deliver telegram', Inches(1), Inches(5.5), Inches(11), Inches(0.6))
add_text(slide, "Tiap 4 jam, sinyal masuk HP. Lo ga perlu buka laptop.", Inches(1), Inches(6.3), Inches(10), Inches(0.5), font_size=16, color=ACCENT)
slide_number(slide, 13)

# ─── SLIDE 14: HANDS-ON ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, "HANDS-ON (15m)", Inches(0.8), Inches(0.5), Inches(10), Inches(0.8), font_size=36, bold=True, color=WHITE)
section_divider(slide, Inches(1.3))
steps = [
    ("Step 1", "Install Hermes", "$ curl -fsSL ... | bash"),
    ("Step 2", "Download template", "QR code / link"),
    ("Step 3", "Jalanin bot", "$ python trading_bot_template.py"),
    ("Step 4", "Cron", "$ hermes cron create ..."),
]
y = Inches(1.8)
for num, title, cmd in steps:
    add_text(slide, num, Inches(1), y, Inches(2), Inches(0.5), font_size=20, bold=True, color=ACCENT)
    add_text(slide, title, Inches(3), y, Inches(4), Inches(0.5), font_size=18, color=WHITE)
    add_text(slide, cmd, Inches(6), y, Inches(6), Inches(0.5), font_size=14, color=GOLD)
    y += Inches(1.0)
add_text(slide, "Selesai! Bot sinyal lo udah jalan sendiri.", Inches(1), Inches(6.0), Inches(10), Inches(0.5), font_size=20, bold=True, color=GOLD)
slide_number(slide, 14)

# ─── SLIDE 15: TIPS ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, "Tips & Best Practices", Inches(0.8), Inches(0.5), Inches(10), Inches(0.8), font_size=36, bold=True, color=WHITE)
section_divider(slide, Inches(1.3))
tips = [
    ("Mulai 1 aset", "XAU atau BTC, jangan semua"),
    ("SL wajib", "Jangan entry tanpa SL"),
    ("Backtest", "Pantau 1 minggu sebelum real trade"),
    ("Cross-check", "AI bisa salah -- cek chart manual"),
    ("Konsisten", "3 sinyal bagus > 20 sinyal random"),
]
y = Inches(1.8)
for title, desc in tips:
    add_text(slide, f"[>] {title}", Inches(1), y, Inches(4), Inches(0.5), font_size=18, bold=True, color=ACCENT)
    add_text(slide, desc, Inches(5), y, Inches(7), Inches(0.5), font_size=16, color=GRAY)
    y += Inches(0.7)
add_text(slide, "Jangan percaya 100% ke AI -- lo tetap bosnya.", Inches(1), Inches(5.5), Inches(10), Inches(0.5), font_size=18, bold=True, color=GOLD)
slide_number(slide, 15)

# ─── SLIDE 16: LIVE BOT GUE ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, "Live Bot Gue: Bukti Nyata", Inches(0.8), Inches(0.5), Inches(10), Inches(0.8), font_size=36, bold=True, color=WHITE)
section_divider(slide, Inches(1.3))
add_text(slide, "Tampilkan dashboard (buka xau_dashboard.html)", Inches(1), Inches(1.8), Inches(11), Inches(0.5), font_size=20, color=GRAY)
add_text(slide, "Ini bot gue - udah generate 5 sinyal, 1 win 4 loss.", Inches(1.2), Inches(2.6), Inches(11), Inches(0.5), font_size=18, color=WHITE)
add_text(slide, "WR 25% tapi profit karena risk management (SL kecil, TP besar).", Inches(1.2), Inches(3.2), Inches(11), Inches(0.5), font_size=18, color=GOLD)
add_text(slide, "Bot tetap jalan walau loss. Yang penting konsisten.", Inches(1.2), Inches(3.8), Inches(11), Inches(0.5), font_size=18, color=WHITE)
add_text(slide, "Link: xau.entry.reasons.html / xau_dashboard.html", Inches(1), Inches(5.0), Inches(10), Inches(0.5), font_size=16, color=ACCENT)
slide_number(slide, 16)

# ─── SLIDE 17: RESOURCES ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, "Resources", Inches(0.8), Inches(0.5), Inches(8), Inches(0.8), font_size=36, bold=True, color=WHITE)
section_divider(slide, Inches(1.3))
resources = [
    ("Hermes Agent", "hermes-agent.nousresearch.com"),
    ("Template Bot", "[Link / QR Code]"),
    ("Cheatsheet Prompt", "[Link / QR Code]"),
    ("PDF GuideBook", "[Link / QR Code]"),
    ("Bot Dashboard", "[Link]"),
]
y = Inches(2.0)
for title, url in resources:
    add_text(slide, title, Inches(1.2), y, Inches(4), Inches(0.5), font_size=18, bold=True, color=ACCENT)
    add_text(slide, url, Inches(5.5), y, Inches(6), Inches(0.5), font_size=16, color=GRAY)
    y += Inches(0.8)
add_text(slide, "Semua link di PDF GuideBook. Scan QR di meja!", Inches(1), Inches(6.0), Inches(10), Inches(0.5), font_size=16, color=GOLD)
slide_number(slide, 17)

# ─── SLIDE 18: Q&A ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
# Big centering
add_text(slide, "Q & A", Inches(0), Inches(2.0), Inches(13.333), Inches(1.5), font_size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "Terima Kasih!", Inches(0), Inches(3.5), Inches(13.333), Inches(1), font_size=36, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(slide, "Lo punya signal generator sendiri sekarang.", Inches(0), Inches(4.8), Inches(13.333), Inches(0.5), font_size=20, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, "Ga perlu jadi programmer. Cukup prompt, download, cron, jalan.", Inches(0), Inches(5.4), Inches(13.333), Inches(0.5), font_size=18, color=GOLD, align=PP_ALIGN.CENTER)
add_text(slide, "@Putra | [Telegram / WA]", Inches(0), Inches(6.3), Inches(13.333), Inches(0.5), font_size=16, color=GRAY, align=PP_ALIGN.CENTER)
slide_number(slide, 18)

# ─── SAVE ───
out_path = os.path.expanduser('~/Desktop/workshop-pintu/FROM-ZERO-TO-TRADING-BOTS.pptx')
prs.save(out_path)
print(f"PPTX saved: {out_path}")
print(f"Size: {os.path.getsize(out_path)} bytes")
print(f"Slides: {len(prs.slides)}")
