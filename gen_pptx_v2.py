#!/usr/bin/env python3
"""Generate PPTX with Hermes-style clean white/blue design."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Palette
WHITE = RGBColor(255, 255, 255)
NEAR_WHITE = RGBColor(250, 251, 253)
BLACK = RGBColor(30, 30, 35)
DARK = RGBColor(40, 45, 55)
BLUE = RGBColor(59, 130, 246)       # primary accent
BLUE_LIGHT = RGBColor(219, 234, 254)
BLUE_DARK = RGBColor(30, 64, 175)
GRAY = RGBColor(120, 125, 140)
GRAY_LIGHT = RGBColor(200, 205, 215)
GREEN = RGBColor(16, 185, 129)
GOLD = RGBColor(245, 158, 11)
BORDER = RGBColor(220, 225, 235)
INNER_BORDER = RGBColor(191, 219, 254)  # light blue

def add_white_card(slide):
    """Full-slide white card with borders."""
    # Outer dark border
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(240, 242, 245)
    shape.line.fill.background()
    # Inner white card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(0.3), Inches(12.733), Inches(6.9))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = BORDER
    card.line.width = Pt(1)
    # Light blue inner border accent
    inner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.38), Inches(0.38), Inches(12.573), Inches(6.74))
    inner.fill.background()
    inner.line.color.rgb = INNER_BORDER
    inner.line.width = Pt(1.5)

def add_section_label(slide, text, left, top):
    add_text(slide, text, left, top, Inches(4), Inches(0.3), font_size=10, color=GRAY, bold=False)

def add_title(slide, text, left, top, width=Inches(10)):
    # Blue accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + Pt(4), Inches(0.04), Inches(0.45))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    add_text(slide, text, left + Inches(0.18), top, width, Inches(0.6), font_size=26, bold=True, color=BLUE_DARK)

def add_text(slide, text, left, top, width, height, font_size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tf

def add_para(tf, text, font_size=16, bold=False, color=BLACK, space_before=Pt(8), font_name='Calibri'):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.space_before = space_before
    return p

def add_bullet(slide, text, left, top, width=Inches(5.5), font_size=15, color=DARK, bold=False):
    tf = add_text(slide, f"  •  {text}", left, top, width, Inches(0.35), font_size=font_size, color=color, bold=bold)
    return tf

def add_step(slide, number, title, desc, left, top, width=Inches(5.5)):
    # Numbered blue circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top + Pt(2), Inches(0.35), Inches(0.35))
    circle.fill.solid()
    circle.fill.fore_color.rgb = BLUE
    circle.line.fill.background()
    tf_c = circle.text_frame
    tf_c.paragraphs[0].text = str(number)
    tf_c.paragraphs[0].font.size = Pt(12)
    tf_c.paragraphs[0].font.bold = True
    tf_c.paragraphs[0].font.color.rgb = WHITE
    tf_c.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf_c.word_wrap = False
    # Title
    add_text(slide, title, left + Inches(0.5), top, width, Inches(0.3), font_size=15, bold=True, color=BLUE_DARK)
    # Desc
    add_text(slide, desc, left + Inches(0.5), top + Inches(0.3), width, Inches(0.6), font_size=13, color=GRAY)

def add_code_block(slide, code_text, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(248, 250, 252)
    shape.line.color.rgb = BLUE_LIGHT
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(14)
    tf.margin_top = Pt(8)
    p = tf.paragraphs[0]
    p.text = code_text
    p.font.size = Pt(13)
    p.font.color.rgb = BLUE_DARK
    p.font.name = 'Consolas'
    return tf

def slide_number(slide, num):
    add_text(slide, f"{num:02d}", Inches(12.1), Inches(6.6), Inches(0.8), Inches(0.4), font_size=11, color=GRAY, align=PP_ALIGN.RIGHT)

def bottom_line(slide, text, left, top):
    add_text(slide, f"— — {text}", left, top, Inches(10), Inches(0.4), font_size=13, color=BLUE)

def separator_line(slide, left, top, width=Inches(11)):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(1))
    line.fill.solid()
    line.fill.fore_color.rgb = GRAY_LIGHT
    line.line.fill.background()

# ═══════════════════════ BUILD SLIDES ═══════════════════════

# ─── SLIDE 1: TITLE ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_card(slide)
add_text(slide, "FROM ZERO", Inches(0.8), Inches(1.5), Inches(10), Inches(1), font_size=54, bold=True, color=BLUE_DARK)
add_text(slide, "TO TRADING BOTS", Inches(0.8), Inches(2.5), Inches(10), Inches(1), font_size=54, bold=True, color=BLUE)
separator_line(slide, Inches(0.8), Inches(3.8))
add_text(slide, "Workshop PINTU x Web3 Dev Bandung", Inches(0.8), Inches(4.1), Inches(10), Inches(0.5), font_size=20, color=BLACK)
add_text(slide, "25 Juli 2026 | 13:00 WIB | Parla Cafe", Inches(0.8), Inches(4.7), Inches(10), Inches(0.4), font_size=16, color=GRAY)
add_text(slide, "Instructor: Putra", Inches(0.8), Inches(5.5), Inches(4), Inches(0.4), font_size=15, bold=True, color=BLUE)
slide_number(slide, 1)

# ─── SLIDE 2: ABOUT ME ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_card(slide)
add_section_label(slide, "ABOUT ME", Inches(0.8), Inches(0.5))
add_title(slide, "Instructor background.", Inches(0.8), Inches(0.9))
# Left column
add_text(slide, "Web Developer & Trader", Inches(0.8), Inches(1.8), Inches(5.5), Inches(0.4), font_size=18, bold=True, color=BLACK)
items = [
    "Full-stack Laravel / PHP sejak 2019",
    "Crypto + jewelry trader",
    "Web3 enthusiast",
    "XAU signal bot live 24/7",
]
y = Inches(2.3)
for item in items:
    add_bullet(slide, item, Inches(0.8), y)
    y += Inches(0.38)
# Right column - bot live stats
card2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.3), Inches(1.8), Inches(5), Inches(2.5))
card2.fill.solid()
card2.fill.fore_color.rgb = RGBColor(239, 246, 255)
card2.line.color.rgb = BLUE_LIGHT
card2.line.width = Pt(1)
add_text(slide, "XAU Bot — LIVE", Inches(7.6), Inches(2.0), Inches(4.5), Inches(0.4), font_size=16, bold=True, color=BLUE_DARK)
stats = [
    "Price: $4,075.70 (+1.41%)",
    "Direction: LONG | Confidence: 95%",
    "Entry: $4,019 | SL: $3,962 | TP: $4,113",
    "R:R 1:1.67 | WR: 25% (1W/4L)",
    "Bot tetap profit walau win rate kecil",
]
y = Inches(2.5)
for s in stats:
    add_text(slide, s, Inches(7.6), y, Inches(4.5), Inches(0.3), font_size=13, color=DARK)
    y += Inches(0.35)
bottom_line(slide, "Bot real, bukan simulasi.", Inches(0.8), Inches(5.8))
slide_number(slide, 2)

# ─── SLIDE 3: AGENDA ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_card(slide)
add_section_label(slide, "AGENDA", Inches(0.8), Inches(0.5))
add_title(slide, "What we'll build in 45 minutes.", Inches(0.8), Inches(0.9))
agenda = [
    ("01", "Install AI Assistant", "Hermes Agent di laptop lo"),
    ("02", "Setup Provider", "OpenRouter API key (gratis)"),
    ("03", "Download Bot Template", "230 baris, 5 indikator, siap jalan"),
    ("04", "Jalanin Bot + Sinyal", "Output sinyal LONG/SHORT + SL/TP"),
    ("05", "Cron — Auto 24/7", "Bot jalan sendiri tanpa server"),
    ("06", "Bonus: Telegram", "Sinyal masuk HP otomatis"),
    ("07", "Q&A", "Tanya apa aja"),
]
for i, (num, title, desc) in enumerate(agenda):
    col = 0 if i < 4 else 1
    row = i if i < 4 else i - 4
    x = Inches(0.8) if col == 0 else Inches(6.8)
    y = Inches(1.7) + row * Inches(0.7)
    add_step(slide, num, title, desc, x, y, width=Inches(5.5))
bottom_line(slide, "45 menit. Dari nol jadi bot jalan.", Inches(0.8), Inches(6.2))
slide_number(slide, 3)

# ─── SLIDE 4: KENAPA BOT ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_card(slide)
add_section_label(slide, "WHY SIGNAL BOT", Inches(0.8), Inches(0.5))
add_title(slide, "Kenapa Bot Sinyal?", Inches(0.8), Inches(0.9))
# Left
add_text(slide, "Lo ga perlu jadi programmer.", Inches(0.8), Inches(1.8), Inches(5.5), Inches(0.4), font_size=18, bold=True, color=BLACK)
items = ["Bisa buka terminal", "Bisa copy-paste", "Mau nyobain"]
y = Inches(2.3)
for item in items:
    add_bullet(slide, item, Inches(0.8), y)
    y += Inches(0.4)
# Right
card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.5), Inches(3.5))
card.fill.solid()
card.fill.fore_color.rgb = RGBColor(239, 246, 255)
card.line.color.rgb = BLUE_LIGHT
card.line.width = Pt(1)
add_text(slide, "Hasil Akhir:", Inches(7.1), Inches(2.0), Inches(5), Inches(0.4), font_size=16, bold=True, color=BLUE_DARK)
result = "Tiap X jam, bot analisa pasar,\ngenerate sinyal, kirim ke Telegram.\nLo tinggal baca, eksekusi kalau cocok."
add_text(slide, result, Inches(7.1), Inches(2.6), Inches(5), Inches(1.5), font_size=15, color=DARK)
add_text(slide, "Ini signal bot, bukan execution bot.", Inches(7.1), Inches(4.2), Inches(5), Inches(0.4), font_size=13, bold=True, color=BLUE)
add_text(slide, "Aman, ga konek API exchange.", Inches(7.1), Inches(4.6), Inches(5), Inches(0.4), font_size=13, color=GRAY)
bottom_line(slide, "AI yang kerja, lo yang megang kontrol.", Inches(0.8), Inches(6.0))
slide_number(slide, 4)

# ─── SLIDE 5: KENAPA HERMES ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_card(slide)
add_section_label(slide, "WHY HERMES AGENT", Inches(0.8), Inches(0.5))
add_title(slide, "Kenapa Hermes Agent?", Inches(0.8), Inches(0.9))
features = [
    ("Open source", "Gratis. Kode terbuka. Ga ada biaya langganan"),
    ("Cron Bawaan", "Bot jalan otomatis tanpa server. Hermes yang jagain"),
    ("Multi-Provider", "Claude, GPT, DeepSeek, Gemini — pilih bebas"),
    ("Multi-Platform", "Telegram, Discord, Email, File — kirim ke mana aja"),
]
for i, (title, desc) in enumerate(features):
    col = 0 if i < 2 else 1
    row = i if i < 2 else i - 2
    x = Inches(0.8) if col == 0 else Inches(6.8)
    y = Inches(1.7) + row * Inches(1.5)
    add_text(slide, title, x, y, Inches(5.5), Inches(0.4), font_size=18, bold=True, color=BLUE_DARK)
    add_text(slide, desc, x, y + Inches(0.4), Inches(5.5), Inches(0.6), font_size=14, color=GRAY)
bottom_line(slide, "Tools yang sama dipake instructor buat bot live-nya.", Inches(0.8), Inches(6.0))
slide_number(slide, 5)

# ─── SLIDE 6: INSTALL ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_card(slide)
add_section_label(slide, "GETTING STARTED", Inches(0.8), Inches(0.5))
add_title(slide, "Install and prep your environment.", Inches(0.8), Inches(0.9))
# Left
add_text(slide, "Kenapa Hermes Agent?", Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.4), font_size=16, bold=True, color=BLACK)
bullets = [
    "Open source — gratis",
    "Cron bawaan — bot jalan otomatis tanpa server",
    "Multi-provider — Claude, GPT, DeepSeek, bebas",
    "Multi-platform — Telegram, Discord, file",
]
y = Inches(2.2)
for b in bullets:
    add_bullet(slide, b, Inches(0.8), y)
    y += Inches(0.38)
add_text(slide, "Install:", Inches(0.8), Inches(4.0), Inches(2), Inches(0.3), font_size=14, bold=True, color=BLACK)
add_code_block(slide, "$ curl -fsSL https://hermes-agent.nousresearch.com/install.sh | bash", Inches(0.8), Inches(4.4), Inches(5.5), Inches(0.55))
add_text(slide, "Cek:", Inches(0.8), Inches(5.1), Inches(2), Inches(0.3), font_size=14, bold=True, color=BLACK)
add_code_block(slide, "$ hermes doctor", Inches(0.8), Inches(5.4), Inches(3.5), Inches(0.5))
# Right - steps
steps = [
    ("Install Python 3.10+", "Download from python.org. Centang 'Add Python to PATH'."),
    ("Set up code editor", "VS Code (free) — ringan, cukup buat edit bot scripts."),
    ("Daftar OpenRouter", "openrouter.ai/keys — gratis. Create Key, copy API key."),
]
for i, (title, desc) in enumerate(steps):
    add_step(slide, i+1, title, desc, Inches(6.8), Inches(1.7) + i * Inches(1.0))
bottom_line(slide, "Once this is done, you're ready to build your first trading bot.", Inches(0.8), Inches(6.2))
slide_number(slide, 6)

# ─── SLIDE 7: SETUP PROVIDER ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_card(slide)
add_section_label(slide, "SETUP PROVIDER", Inches(0.8), Inches(0.5))
add_title(slide, "Connect Hermes to an AI model.", Inches(0.8), Inches(0.9))
# Left
add_text(slide, "Step 1: OpenRouter (GRATIS)", Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.4), font_size=16, bold=True, color=BLUE_DARK)
bullets = [
    "Buka https://openrouter.ai/keys",
    "Login (Google/GitHub/email)",
    "Klik 'Create Key' → Copy API key",
    "Simpan di Notepad",
]
y = Inches(2.2)
for b in bullets:
    add_bullet(slide, b, Inches(0.8), y)
    y += Inches(0.38)
add_text(slide, "Step 2: Setup di Hermes", Inches(0.8), Inches(4.0), Inches(5.5), Inches(0.4), font_size=16, bold=True, color=BLUE_DARK)
add_code_block(slide, "$ hermes model", Inches(0.8), Inches(4.5), Inches(5.5), Inches(0.55))
add_text(slide, "Pilih: OpenRouter -> paste API key -> pilih model", Inches(1), Inches(5.2), Inches(5.5), Inches(0.4), font_size=14, color=DARK)
# Right
card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.7), Inches(5.5), Inches(4.0))
card.fill.solid()
card.fill.fore_color.rgb = RGBColor(239, 246, 255)
card.line.color.rgb = BLUE_LIGHT
card.line.width = Pt(1)
add_text(slide, "Model Gratis Rekomendasi", Inches(7.1), Inches(1.9), Inches(5), Inches(0.4), font_size=16, bold=True, color=BLUE_DARK)
models = [
    "DeepSeek V3  →  signal analisa (recommended!)",
    "Qwen 2.5 72B  →  coding sederhana",
    "Gemini Flash  →  gratis unlimited",
]
y = Inches(2.5)
for m in models:
    add_text(slide, m, Inches(7.1), y, Inches(5), Inches(0.35), font_size=14, color=DARK)
    y += Inches(0.45)
separator_line(slide, Inches(7.1), Inches(4.0))
add_text(slide, "Alternatif:", Inches(7.1), Inches(4.2), Inches(5), Inches(0.3), font_size=14, bold=True, color=BLACK)
add_text(slide, "Gemini API — 60 req/menit gratis", Inches(7.1), Inches(4.6), Inches(5), Inches(0.3), font_size=13, color=GRAY)
add_text(slide, "Ollama local — butuh RAM 8GB+", Inches(7.1), Inches(5.0), Inches(5), Inches(0.3), font_size=13, color=GRAY)
bottom_line(slide, "DeepSeek V3 — gratis, cepet, cukup buat workshop.", Inches(0.8), Inches(6.2))
slide_number(slide, 7)

# ─── SLIDE 8: TEMPLATE ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_card(slide)
add_section_label(slide, "BOT TEMPLATE", Inches(0.8), Inches(0.5))
add_title(slide, "Template bot — 230 baris, siap jalan.", Inches(0.8), Inches(0.9))
# Left
add_text(slide, "trading_bot_template.py", Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.4), font_size=16, bold=True, color=BLUE_DARK)
items = [
    "Fetch harga (demo mode, tinggal ganti API real)",
    "5 Indikator: RSI, EMA, MACD, BB, ATR",
    "Generate sinyal LONG / SHORT + SL / TP",
    "Simpan history ke signals.json",
    "Cooldown 4 jam + silent mode (ga spam)",
]
y = Inches(2.2)
for item in items:
    add_bullet(slide, item, Inches(0.8), y, width=Inches(5.5))
    y += Inches(0.38)
add_text(slide, "Ga perlu nulis kode. Tinggal download dan jalanin.", Inches(0.8), Inches(4.5), Inches(5.5), Inches(0.4), font_size=15, bold=True, color=BLUE)
# Right - code snippet
add_code_block(slide, (
    "INDICATORS = {\n"
    '    "RSI": "Relative Strength Index",\n'
    '    "EMA": "Exponential Moving Average",\n'
    '    "MACD": "MACD momentum",\n'
    '    "BB": "Bollinger Bands envelope",\n'
    '    "ATR": "Average True Range",\n'
    "}"
), Inches(6.8), Inches(1.7), Inches(5.5), Inches(2.0))
add_code_block(slide, (
    "def generate_signal(prices, history, last_signal):\n"
    '    """Combine indikator -> confidence 0-100%.\n'
    "    >=60% -> LONG/SHORT signal.\n"
    '    <60% -> skip."""'
), Inches(6.8), Inches(4.0), Inches(5.5), Inches(1.2))
bottom_line(slide, "Template ini yang bakal lo jalanin + cron-in.", Inches(0.8), Inches(6.2))
slide_number(slide, 8)

# ─── SLIDE 9: JALANIN BOT ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_card(slide)
add_section_label(slide, "RUN BOT", Inches(0.8), Inches(0.5))
add_title(slide, "Jalanin bot dan liat sinyal.", Inches(0.8), Inches(0.9))
add_code_block(slide, "$ python trading_bot_template.py", Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.55))
# Output card
out_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.7), Inches(5.5), Inches(2.6))
out_card.fill.solid()
out_card.fill.fore_color.rgb = RGBColor(239, 246, 255)
out_card.line.color.rgb = BLUE_LIGHT
out_card.line.width = Pt(1)
add_text(slide, "Output:", Inches(7.1), Inches(1.9), Inches(5), Inches(0.3), font_size=13, bold=True, color=BLUE_DARK)
output_text = (
    "TRADING SIGNAL - 25 Jul 13:30 WIB\n"
    "========================================\n"
    "Price: $4,075.70 (+1.71%)\n"
    "Direction: LONG | Confidence: 95%\n"
    "Entry: $4,075.70 | SL: $3,962.50\n"
    "TP: $4,113.20 | R:R 1:1.67"
)
add_text(slide, output_text, Inches(7.1), Inches(2.3), Inches(5), Inches(1.8), font_size=13, color=BLUE_DARK, font_name='Consolas')
add_text(slide, "Bot juga simpan history di:", Inches(0.8), Inches(4.5), Inches(5.5), Inches(0.3), font_size=14, color=GRAY)
add_text(slide, "signals.json + history.json", Inches(0.8), Inches(4.9), Inches(5.5), Inches(0.3), font_size=14, bold=True, color=BLUE_DARK, font_name='Consolas')
bottom_line(slide, "Kalau bot silent — itu normal (cooldown).", Inches(0.8), Inches(6.2))
slide_number(slide, 9)

# ─── SLIDE 10: BOT SILENT ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_card(slide)
add_section_label(slide, "BOT BEHAVIOR", Inches(0.8), Inches(0.5))
add_title(slide, "Kenapa bot ga ngomong tiap detik?", Inches(0.8), Inches(0.9))
add_text(slide, "Itu FITUR, bukan bug.", Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.4), font_size=20, bold=True, color=GREEN)
# Cards
reasons = [
    ("Cooldown 4 jam", "Jeda antar sinyal biar ga spam. Bot ga akan generate sinyal baru sebelum 4 jam."),
    ("Active guard", "Ga generate sinyal kalo masih ada posisi aktif yang belum kena TP/SL."),
    ("Silent mode", "Output cuma kalo ada sinyal BARU atau posisi CLOSED. Selain itu diem aja."),
]
for i, (title, desc) in enumerate(reasons):
    x = Inches(0.8) + (i % 3) * Inches(4.1)
    y = Inches(2.4)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.8), Inches(2.0))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(239, 246, 255)
    card.line.color.rgb = BLUE_LIGHT
    card.line.width = Pt(1)
    add_text(slide, title, x + Inches(0.2), y + Inches(0.15), Inches(3.4), Inches(0.4), font_size=15, bold=True, color=BLUE_DARK)
    add_text(slide, desc, x + Inches(0.2), y + Inches(0.6), Inches(3.4), Inches(1.2), font_size=13, color=DARK)
add_text(slide, "Mau liat status? Buka signals.json atau jalanin ulang bot.", Inches(0.8), Inches(5.0), Inches(10), Inches(0.4), font_size=15, color=GRAY)
slide_number(slide, 10)

# ─── SLIDE 11: LIVE DEMO ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_card(slide)
add_section_label(slide, "LIVE DEMO", Inches(0.8), Inches(0.5))
add_title(slide, "Bot XAU — hidup 24/7.", Inches(0.8), Inches(0.9))
# Left - bot card
card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.5), Inches(3.2))
card.fill.solid()
card.fill.fore_color.rgb = RGBColor(239, 246, 255)
card.line.color.rgb = BLUE_LIGHT
card.line.width = Pt(1)
add_text(slide, "XAU SIGNAL — ACTIVE", Inches(1.1), Inches(2.0), Inches(5), Inches(0.4), font_size=16, bold=True, color=BLUE_DARK)
live_data = "LONG | Entry $4,019\nP&L +1.41%\nSL $3,962 | TP $4,113\n36/48h | R:R 1:1.67"
add_text(slide, live_data, Inches(1.1), Inches(2.6), Inches(5), Inches(1.5), font_size=15, color=DARK, font_name='Consolas')
add_text(slide, "WR: 25% (1W/4L)", Inches(1.1), Inches(4.2), Inches(5), Inches(0.4), font_size=14, bold=True, color=GOLD)
# Right
add_text(slide, "Kenapa lo harus percaya?", Inches(6.8), Inches(1.8), Inches(5.5), Inches(0.4), font_size=16, bold=True, color=BLUE_DARK)
reasons = [
    "Ini bot real, bukan screenshot",
    "Gue jalanin sebelum workshop",
    "5 sinyal udah di-generate",
    "1 win, 4 loss — tapi profit karena SL kecil, TP besar",
    "Bot tetap jalan walau loss",
]
y = Inches(2.4)
for r in reasons:
    add_bullet(slide, r, Inches(6.8), y, width=Inches(5.5))
    y += Inches(0.4)
add_text(slide, "Gue ga pernah coding ulang. Template -> Cron -> Done.", Inches(6.8), Inches(5.0), Inches(5.5), Inches(0.4), font_size=15, bold=True, color=BLUE)
slide_number(slide, 11)

# ─── SLIDE 12: CRON ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_card(slide)
add_section_label(slide, "AUTOMATION", Inches(0.8), Inches(0.5))
add_title(slide, "Cron: bot jalan sendiri, lo tinggal terima hasil.", Inches(0.8), Inches(0.9))
add_text(slide, "Manual python trading_bot.py tiap hari? Capek.", Inches(0.8), Inches(1.7), Inches(10), Inches(0.4), font_size=16, color=GRAY)
add_text(slide, "Cron = bot jalan sendiri tanpa server.", Inches(0.8), Inches(2.2), Inches(10), Inches(0.4), font_size=18, bold=True, color=BLUE_DARK)
add_code_block(slide, (
    "$ hermes cron create 'every 4h' \\\n"
    '    --name "Sinyal XAU Auto" \\\n'
    '    --prompt "Jalankan python trading_bot_template.py 2>&1. Kirim output apa adanya."'
), Inches(0.8), Inches(2.8), Inches(11), Inches(1.5))
add_text(slide, "Cek cron yang berjalan:", Inches(0.8), Inches(4.6), Inches(4), Inches(0.4), font_size=14, color=GRAY)
add_code_block(slide, "$ hermes cron list", Inches(0.8), Inches(5.0), Inches(4.5), Inches(0.5))
bottom_line(slide, "Cron -> sinyal otomatis. Lo tinggal baca hasilnya.", Inches(0.8), Inches(6.2))
slide_number(slide, 12)

# ─── SLIDE 13: TELEGRAM ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_card(slide)
add_section_label(slide, "NOTIFICATIONS", Inches(0.8), Inches(0.5))
add_title(slide, "Kirim sinyal ke Telegram.", Inches(0.8), Inches(0.9))
add_text(slide, "Setup Gateway:", Inches(0.8), Inches(1.7), Inches(5), Inches(0.4), font_size=16, bold=True, color=BLUE_DARK)
steps = [
    ("Bikin bot lewat @BotFather", "Dapetin token"),
    ("Start bot di HP lo", "Cari nama bot lo, klik start"),
    ("Setup Hermes Gateway", "$ hermes gateway setup -> pilih Telegram"),
    ("Update cron biar deliver", "$ hermes cron update <id> --deliver telegram"),
]
for i, (title, desc) in enumerate(steps):
    add_step(slide, i+1, title, desc, Inches(0.8), Inches(2.2) + i * Inches(0.65), width=Inches(5.5))
# Right - visual
card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.3), Inches(1.7), Inches(5), Inches(3.5))
card.fill.solid()
card.fill.fore_color.rgb = RGBColor(239, 246, 255)
card.line.color.rgb = BLUE_LIGHT
card.line.width = Pt(1)
add_text(slide, "Hasilnya:", Inches(7.6), Inches(1.9), Inches(4.5), Inches(0.4), font_size=16, bold=True, color=BLUE_DARK)
msg = "[TRADING SIGNAL - 25 Jul 14:00 WIB]\n"
msg += "Price: $4,075.70 (+1.71%)\n"
msg += "Direction: LONG | Confidence: 95%\n"
msg += "Entry: $4,075.70 | R:R 1:1.67"
add_text(slide, msg, Inches(7.6), Inches(2.5), Inches(4.5), Inches(1.5), font_size=13, color=DARK, font_name='Consolas')
add_text(slide, "Tiap 4 jam, sinyal masuk HP.", Inches(7.6), Inches(4.3), Inches(4.5), Inches(0.4), font_size=16, bold=True, color=BLUE)
add_text(slide, "Lo ga perlu buka laptop.", Inches(7.6), Inches(4.7), Inches(4.5), Inches(0.4), font_size=14, color=GRAY)
slide_number(slide, 13)

# ─── SLIDE 14: HANDS-ON ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_card(slide)
add_section_label(slide, "HANDS-ON", Inches(0.8), Inches(0.5))
add_title(slide, "Practice time — 15 menit.", Inches(0.8), Inches(0.9))
steps = [
    ("Install Hermes", "curl -fsSL https://hermes-agent.nousresearch.com/install.sh | bash"),
    ("Download template", "QR code / link workshop"),
    ("Jalanin bot", "python trading_bot_template.py"),
    ("Setup cron", "hermes cron create 'every 6h' --name 'Sinyal Saya' --prompt ..."),
    ("Telegram (opsional)", "hermes gateway setup -> deliver ke Telegram"),
]
for i, (title, cmd) in enumerate(steps):
    col = 0 if i < 3 else 1
    row = i if i < 3 else i - 3
    x = Inches(0.8) if col == 0 else Inches(6.8)
    y = Inches(1.7) + row * Inches(1.1)
    add_step(slide, i+1, title, cmd, x, y, width=Inches(5.5))
add_text(slide, "Selesai! Bot sinyal lo udah jalan sendiri.", Inches(0.8), Inches(5.8), Inches(10), Inches(0.5), font_size=20, bold=True, color=BLUE)
slide_number(slide, 14)

# ─── SLIDE 15: TIPS ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_card(slide)
add_section_label(slide, "BEST PRACTICES", Inches(0.8), Inches(0.5))
add_title(slide, "Tips biar trading bot lo sukses.", Inches(0.8), Inches(0.9))
# Left - do
add_text(slide, "DO:", Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.4), font_size=18, bold=True, color=GREEN)
do_items = [
    "Mulai dengan 1 aset (XAU atau BTC)",
    "SL wajib sebelum entry",
    "Backtest 1 minggu sebelum real trade",
    "Cross-check sinyal AI dengan chart manual",
    "Konsisten — 3 sinyal bagus > 20 random",
]
y = Inches(2.2)
for item in do_items:
    add_bullet(slide, item, Inches(0.8), y, width=Inches(5.5))
    y += Inches(0.38)
# Right - dont
add_text(slide, "DON'T:", Inches(6.8), Inches(1.7), Inches(5.5), Inches(0.4), font_size=18, bold=True, color=RGBColor(220, 38, 38))
dont_items = [
    "Trade real money langsung",
    "Over-trading (cron tiap jam = spam)",
    "Percaya 100% ke AI — lo tetap bosnya",
]
y = Inches(2.2)
for item in dont_items:
    add_text(slide, f"  X  {item}", Inches(6.8), y, Inches(5.5), Inches(0.35), font_size=15, color=RGBColor(200, 50, 50))
    y += Inches(0.4)
bottom_line(slide, "Risk management > win rate.", Inches(0.8), Inches(6.2))
slide_number(slide, 15)

# ─── SLIDE 16: BUKTI ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_card(slide)
add_section_label(slide, "REAL PROOF", Inches(0.8), Inches(0.5))
add_title(slide, "Bot gue — bukti nyata.", Inches(0.8), Inches(0.9))
# Dashboard
add_text(slide, "Tampilkan dashboard:", Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.4), font_size=16, bold=True, color=BLUE_DARK)
add_text(slide, "xau_dashboard.html", Inches(1), Inches(2.2), Inches(5), Inches(0.3), font_size=15, color=BLUE, font_name='Consolas')
add_text(slide, "xau_entry_reasons.html", Inches(1), Inches(2.6), Inches(5), Inches(0.3), font_size=15, color=BLUE, font_name='Consolas')
# Stats
card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.7), Inches(5.5), Inches(2.8))
card.fill.solid()
card.fill.fore_color.rgb = RGBColor(239, 246, 255)
card.line.color.rgb = BLUE_LIGHT
card.line.width = Pt(1)
add_text(slide, "Bot Stats", Inches(7.1), Inches(1.9), Inches(5), Inches(0.4), font_size=16, bold=True, color=BLUE_DARK)
stats = [
    "Total sinyal: 5",
    "Win: 1 | Loss: 4",
    "Win Rate: 25%",
    "Tapi profit karena R:R bagus",
    "Bot tetap jalan walau loss",
]
y = Inches(2.5)
for s in stats:
    add_text(slide, f"  >  {s}", Inches(7.1), y, Inches(5), Inches(0.35), font_size=14, color=DARK)
    y += Inches(0.38)
add_text(slide, "Yang penting konsisten. Bukan selalu benar.", Inches(1), Inches(5.0), Inches(10), Inches(0.4), font_size=16, bold=True, color=BLUE)
slide_number(slide, 16)

# ─── SLIDE 17: RESOURCES ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_card(slide)
add_section_label(slide, "RESOURCES", Inches(0.8), Inches(0.5))
add_title(slide, "Semua yang lo butuh setelah workshop.", Inches(0.8), Inches(0.9))
resources = [
    ("Hermes Agent", "hermes-agent.nousresearch.com"),
    ("OpenRouter", "openrouter.ai/keys"),
    ("Python Download", "python.org/downloads"),
    ("Bot Template", "[QR Code / Link]"),
    ("Cheatsheet Prompt", "[QR Code / Link]"),
    ("PDF GuideBook", "[QR Code / Link]"),
]
cols = 2
for i, (title, url) in enumerate(resources):
    col = i % cols
    row = i // cols
    x = Inches(0.8) if col == 0 else Inches(6.8)
    y = Inches(1.7) + row * Inches(0.9)
    add_text(slide, title, x, y, Inches(5.5), Inches(0.35), font_size=16, bold=True, color=BLUE_DARK)
    add_text(slide, url, x, y + Inches(0.35), Inches(5.5), Inches(0.35), font_size=14, color=BLUE)
add_text(slide, "Semua link ada di PDF GuideBook. Scan QR!", Inches(0.8), Inches(5.5), Inches(10), Inches(0.4), font_size=15, bold=True, color=GOLD)
slide_number(slide, 17)

# ─── SLIDE 18: Q&A ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_white_card(slide)
add_text(slide, "Q & A", Inches(0), Inches(2.2), Inches(13.333), Inches(1), font_size=56, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
add_text(slide, "Terima Kasih!", Inches(0), Inches(3.4), Inches(13.333), Inches(0.7), font_size=32, color=BLUE_DARK, align=PP_ALIGN.CENTER)
separator_line(slide, Inches(4), Inches(4.3), Inches(5.333))
add_text(slide, "Lo punya signal generator sendiri sekarang.", Inches(0), Inches(4.6), Inches(13.333), Inches(0.5), font_size=18, color=DARK, align=PP_ALIGN.CENTER)
add_text(slide, "Ga perlu jadi programmer. Cukup prompt, download, cron, jalan.", Inches(0), Inches(5.2), Inches(13.333), Inches(0.5), font_size=16, color=BLUE, align=PP_ALIGN.CENTER)
add_text(slide, "@Putra", Inches(0), Inches(6.0), Inches(13.333), Inches(0.4), font_size=15, color=GRAY, align=PP_ALIGN.CENTER)
slide_number(slide, 18)

# ─── SAVE ───
out_path = os.path.expanduser('~/Desktop/workshop-pintu/FROM-ZERO-TO-TRADING-BOTS.pptx')
prs.save(out_path)
print(f"PPTX saved: {out_path}")
print(f"Size: {os.path.getsize(out_path)} bytes")
print(f"Slides: {len(prs.slides)}")
