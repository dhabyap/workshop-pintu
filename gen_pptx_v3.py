#!/usr/bin/env python3
"""Generate PPTX — persis style PDF: putih, aksen biru-ungu, clean sans-serif."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Palette
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
DARK = RGBColor(30, 30, 35)
GRAY = RGBColor(107, 114, 128)
GRAY_LIGHT = RGBColor(160, 165, 175)
LINE_COLOR = RGBColor(224, 224, 230)
ACCENT = RGBColor(99, 102, 241)     # biru-ungu (#6366F1)
ACCENT_LIGHT = RGBColor(199, 210, 254)
ACCENT_DARK = RGBColor(67, 56, 202)
GREEN = RGBColor(16, 185, 129)
RED = RGBColor(239, 68, 68)
GOLD = RGBColor(245, 158, 11)
CODE_BG = RGBColor(248, 249, 252)

def add_text(slide, text, left, top, width, height, font_size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT, font_name='Calibri', italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.font.italic = italic
    p.alignment = align
    return tf

def add_para(tf, text, font_size=16, bold=False, color=DARK, space_before=Pt(6), font_name='Calibri', italic=False):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.font.italic = italic
    p.space_before = space_before
    return p

def add_section_label(slide, text, left, top):
    """Small uppercase category label like 'GETTING STARTED'"""
    add_text(slide, text.upper(), left, top, Inches(6), Inches(0.35), font_size=11, color=ACCENT, bold=True)

def add_title(slide, text, left, top, width=Inches(10)):
    """Big bold title with vertical accent bar on left"""
    # Vertical accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + Pt(2), Pt(4), Inches(0.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    # Title text
    add_text(slide, text, left + Inches(0.2), top, width, Inches(0.7), font_size=30, bold=True, color=BLACK)

def add_subtitle(slide, text, left, top, width=Inches(10)):
    """Gray subtitle below title"""
    add_text(slide, text, left, top, width, Inches(0.4), font_size=16, color=GRAY)

def add_separator(slide, left, top, width=Inches(11.5)):
    """Thin light gray horizontal line"""
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(1))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE_COLOR
    line.line.fill.background()

def add_footer(slide, text, left, top, width=Inches(11)):
    """Gray footer text after separator"""
    add_text(slide, text, left, top, width, Inches(0.4), font_size=11, color=GRAY, italic=True)

def add_check_item(slide, text, left, top, width=Inches(10), font_size=16, bold=False, color=DARK):
    """Checkmark-in-circle icon + text"""
    # Blue-purple circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top + Pt(3), Inches(0.28), Inches(0.28))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT
    circle.line.fill.background()
    # Checkmark inside
    tf_c = circle.text_frame
    tf_c.paragraphs[0].text = ' '
    tf_c.paragraphs[0].font.size = Pt(1)
    tf_c.word_wrap = False
    # Text
    add_text(slide, text, left + Inches(0.45), top, width, Inches(0.35), font_size=font_size, color=color, bold=bold)

def add_step_item(slide, num, title, desc, left, top, width=Inches(5)):
    """Numbered step: bold number + bold title + desc below"""
    # Number
    add_text(slide, f"0{num}", left, top, Inches(0.6), Inches(0.35), font_size=14, bold=True, color=ACCENT)
    # Title
    add_text(slide, title, left + Inches(0.7), top, width, Inches(0.25), font_size=15, bold=True, color=BLACK)
    # Desc
    add_text(slide, desc, left + Inches(0.7), top + Inches(0.28), width, Inches(0.4), font_size=12, color=GRAY)

def add_code_block(slide, code_text, left, top, width, height):
    """Code block with light gray background"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(14)
    tf.margin_top = Pt(6)
    p = tf.paragraphs[0]
    p.text = code_text
    p.font.size = Pt(11)
    p.font.color.rgb = DARK
    p.font.name = 'Consolas'
    return tf

def add_card(slide, text, left, top, width, height, color=ACCENT):
    """Rounded card with color fill"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_top = Pt(8)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE
    p.font.name = 'Calibri'
    p.font.bold = True
    return tf

def slide_number(slide, num, total=18):
    """Bottom-right page number"""
    add_text(slide, f"{num} / {total}", Inches(12.2), Inches(6.9), Inches(1), Inches(0.3), font_size=9, color=GRAY, align=PP_ALIGN.RIGHT)

# ═══════════════════════ BUILD SLIDES ═══════════════════════

# ─── SLIDE 1: TITLE (cover) ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
# White background
add_text(slide, "WORKSHOP · PINTU × WEB3 DEV BANDUNG", Inches(0.8), Inches(1.2), Inches(8), Inches(0.5), font_size=14, color=ACCENT, bold=True)
add_title(slide, "From Zero to Trading Bots", Inches(0.8), Inches(1.9), Inches(10))
add_subtitle(slide, "Build your first crypto trading bot today — no prior experience needed!", Inches(1.0), Inches(2.7))
add_text(slide, "Hosted by PINTU × Web3 Dev Bandung", Inches(0.8), Inches(4.2), Inches(8), Inches(0.4), font_size=16, bold=True, color=BLACK)
add_text(slide, "Beginner-friendly · Hands-on · Live Coding Session", Inches(0.8), Inches(4.7), Inches(8), Inches(0.4), font_size=13, color=GRAY)
add_text(slide, "Sabtu, 25 Juli 2026 | 13:00 WIB | Parla Cafe", Inches(0.8), Inches(5.5), Inches(6), Inches(0.4), font_size=14, color=DARK)
add_text(slide, "Instructor: Putra", Inches(0.8), Inches(6.0), Inches(4), Inches(0.4), font_size=13, color=ACCENT, bold=True)
slide_number(slide, 1)

# ─── SLIDE 2: ABOUT ME ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_label(slide, "GETTING STARTED", Inches(0.8), Inches(0.5))
add_title(slide, "About Me", Inches(0.8), Inches(0.95))
add_subtitle(slide, "Today I'll show you how to build your own AI trading signal bot from scratch — no complex AI knowledge required.", Inches(1.0), Inches(1.8))
add_separator(slide, Inches(0.8), Inches(2.5))
# Bullet items
items = [
    "Web developer sejak 2019 (Laravel, PHP, Web3)",
    "Crypto + jewelry trader — manage both portfolios",
    "Run XAU signal bot live 24/7 — real signals, real P&L",
    "5 indicators: RSI, EMA, MACD, Bollinger Bands, ATR",
    "Bot sends signals to Telegram — hands-free monitoring",
]
y = Inches(2.8)
for item in items:
    add_check_item(slide, item, Inches(0.8), y)
    y += Inches(0.5)
add_separator(slide, Inches(0.8), Inches(5.3))
add_footer(slide, "All source code and setup guides will be shared after the session.", Inches(0.8), Inches(5.6))
slide_number(slide, 2)

# ─── SLIDE 3: WHAT YOU WILL BUILD ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_label(slide, "WHAT YOU WILL BUILD TODAY", Inches(0.8), Inches(0.5))
add_title(slide, "By the end of this workshop, you will have built:", Inches(0.8), Inches(0.95))
items = [
    "An AI-powered trading signal bot that analyzes markets",
    "Custom indicators: RSI, EMA, MACD, Bollinger Bands, ATR",
    "Automated cron job — bot runs every 4 hours without server",
    "Optional Telegram integration — signals delivered to your phone",
    "A complete understanding of how to modify and improve your bot",
]
y = Inches(1.8)
for item in items:
    add_check_item(slide, item, Inches(0.8), y, font_size=15)
    y += Inches(0.7)
add_separator(slide, Inches(0.8), Inches(5.5))
add_footer(slide, "All source code and setup guides will be shared with you via the workshop group after the session.", Inches(0.8), Inches(5.8))
slide_number(slide, 3)

# ─── SLIDE 4: AGENDA ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_label(slide, "TODAY'S PLAN", Inches(0.8), Inches(0.5))
add_title(slide, "How this workshop is structured.", Inches(0.8), Inches(0.95))
agenda = [
    ("01", "Check-in & Setup Prep", "5 min"),
    ("02", "Install AI Assistant (Hermes Agent)", "5 min"),
    ("03", "Connect to AI Provider (OpenRouter)", "5 min"),
    ("04", "Run Your First Bot Signal", "10 min"),
    ("05", "Cron: Automate the Bot", "10 min"),
    ("06", "Bonus: Telegram Integration", "5 min"),
    ("07", "Q&A", "5 min"),
]
y = Inches(1.7)
for num, title, dur in agenda:
    # Line separator
    add_separator(slide, Inches(0.8), y, Inches(11))
    add_text(slide, num, Inches(0.8), y + Pt(4), Inches(0.6), Inches(0.35), font_size=14, bold=True, color=ACCENT)
    add_text(slide, title, Inches(1.6), y + Pt(4), Inches(7), Inches(0.35), font_size=15, color=DARK)
    add_text(slide, dur, Inches(10.5), y + Pt(4), Inches(1.5), Inches(0.35), font_size=14, color=GRAY, align=PP_ALIGN.RIGHT)
    y += Inches(0.55)
add_text(slide, "Total session time: 45 minutes", Inches(0.8), Inches(5.8), Inches(5), Inches(0.4), font_size=12, color=GRAY)
add_text(slide, "Timing is approximate. Feel free to ask questions throughout.", Inches(6), Inches(5.8), Inches(6), Inches(0.4), font_size=11, color=GRAY, align=PP_ALIGN.RIGHT)
slide_number(slide, 4)

# ─── SLIDE 5: WHY SIGNAL BOT ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_label(slide, "WHY SIGNAL BOT", Inches(0.8), Inches(0.5))
add_title(slide, "Kenapa bot sinyal — dan kenapa lo bisa buatnya?", Inches(0.8), Inches(0.95))
# Left column
add_text(slide, "No coding experience needed. Just:", Inches(0.8), Inches(1.8), Inches(5.5), Inches(0.4), font_size=16, bold=True, color=DARK)
items = [
    "Open a terminal",
    "Copy-paste commands",
    "Follow along",
]
y = Inches(2.3)
for item in items:
    add_check_item(slide, item, Inches(0.8), y, font_size=15)
    y += Inches(0.45)
# Right column
card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.7), Inches(5.5), Inches(3.0))
card.fill.solid()
card.fill.fore_color.rgb = RGBColor(238, 242, 255)
card.line.fill.background()
add_text(slide, "What you'll get:", Inches(7.1), Inches(1.9), Inches(5), Inches(0.4), font_size=16, bold=True, color=ACCENT)
results = [
    "Signal bot analyzes market every 4h",
    "Generates LONG/SHORT signals with SL/TP",
    "Sends to Terminal (and optionally Telegram)",
    "You decide — execute or skip",
    "Zero risk: signal-only, no exchange connection",
]
y = Inches(2.4)
for r in results:
    add_check_item(slide, r, Inches(7.1), y, width=Inches(5), font_size=14)
    y += Inches(0.42)
add_separator(slide, Inches(7.1), y + Inches(0.15), Inches(4.5))
add_text(slide, "AI generates signals. You stay in control.", Inches(7.1), y + Inches(0.3), Inches(5), Inches(0.4), font_size=12, color=ACCENT, italic=True)
add_separator(slide, Inches(0.8), Inches(5.5))
add_footer(slide, "Signal bot = analysis only. No API keys for exchanges needed.", Inches(0.8), Inches(5.8))
slide_number(slide, 5)

# ─── SLIDE 6: WHY HERMES ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_label(slide, "TOOLS", Inches(0.8), Inches(0.5))
add_title(slide, "Kenapa Hermes Agent?", Inches(0.8), Inches(0.95))
features = [
    ("Open Source", "Gratis, kode terbuka, ga ada biaya langganan"),
    ("Built-in Cron", "Bot jalan otomatis tanpa server — Hermes jagain"),
    ("Multi-Provider", "DeepSeek, Claude, GPT, Gemini — switch anytime"),
    ("Multi-Platform", "Telegram, Discord, Email, File — deliver anywhere"),
]
for i, (title, desc) in enumerate(features):
    col = 0 if i < 2 else 1
    row = i if i < 2 else i - 2
    x = Inches(0.8) if col == 0 else Inches(6.8)
    y = Inches(1.7) + row * Inches(1.3)
    add_card(slide, title, x, y, Inches(5.5), Inches(0.45), color=ACCENT)
    add_text(slide, desc, x, y + Inches(0.55), Inches(5.5), Inches(0.4), font_size=13, color=GRAY)
add_separator(slide, Inches(0.8), Inches(5.5))
add_footer(slide, "Same tool used by the instructor for his live XAU bot.", Inches(0.8), Inches(5.8))
slide_number(slide, 6)

# ─── SLIDE 7: INSTALL ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_label(slide, "GETTING STARTED", Inches(0.8), Inches(0.5))
add_title(slide, "Install and prep your environment.", Inches(0.8), Inches(0.95))
# Left
add_text(slide, "Install Hermes Agent:", Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.4), font_size=15, bold=True, color=DARK)
add_code_block(slide, "curl -fsSL https://hermes-agent.nousresearch.com/install.sh | bash", Inches(0.8), Inches(2.2), Inches(5.5), Inches(0.5))
add_text(slide, "Verify:", Inches(0.8), Inches(2.9), Inches(2), Inches(0.3), font_size=14, bold=True, color=DARK)
add_code_block(slide, "hermes doctor", Inches(0.8), Inches(3.2), Inches(3.5), Inches(0.4))
add_text(slide, "Windows: Python 3.10+ required.", Inches(0.8), Inches(3.8), Inches(5.5), Inches(0.3), font_size=12, color=GRAY)
# Right - steps
steps = [
    ("Install Python 3.10+", "Download from python.org. Check 'Add to PATH'."),
    ("Set up VS Code", "Free, lightweight. Install if you want to edit scripts."),
    ("Get OpenRouter API key", "openrouter.ai/keys — free. Create account, copy key."),
]
for i, (title, desc) in enumerate(steps):
    add_step_item(slide, i+1, title, desc, Inches(6.8), Inches(1.7) + i * Inches(0.65))
add_separator(slide, Inches(0.8), Inches(5.5))
add_footer(slide, "Once this is done, you're ready to build your first trading bot.", Inches(0.8), Inches(5.8))
slide_number(slide, 7)

# ─── SLIDE 8: SETUP PROVIDER ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_label(slide, "SETUP PROVIDER", Inches(0.8), Inches(0.5))
add_title(slide, "Connect Hermes to an AI model.", Inches(0.8), Inches(0.95))
# Left
add_text(slide, "OpenRouter (GRATIS):", Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.4), font_size=15, bold=True, color=DARK)
items = [
    "Open https://openrouter.ai/keys",
    "Login with Google/GitHub/email",
    "Click 'Create Key' → Copy the API key",
    "Save it — you'll paste it into Hermes",
]
y = Inches(2.2)
for item in items:
    add_check_item(slide, item, Inches(0.8), y, font_size=14)
    y += Inches(0.4)
add_code_block(slide, "hermes model", Inches(0.8), Inches(4.0), Inches(5.5), Inches(0.4))
add_text(slide, "Select: OpenRouter → paste key → pick DeepSeek V3", Inches(1.0), Inches(4.5), Inches(5.5), Inches(0.4), font_size=13, color=DARK)
# Right
card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.7), Inches(5.5), Inches(3.2))
card.fill.solid()
card.fill.fore_color.rgb = RGBColor(238, 242, 255)
card.line.fill.background()
add_text(slide, "Recommended Free Models:", Inches(7.1), Inches(1.9), Inches(5), Inches(0.4), font_size=15, bold=True, color=ACCENT)
models = [
    ("DeepSeek V3", "Best for signal analysis"),
    ("Qwen 2.5 72B", "Good for code modifications"),
    ("Gemini Flash", "Fast, unlimited free tier"),
]
y = Inches(2.5)
for title, desc in models:
    add_text(slide, title, Inches(7.1), y, Inches(5), Inches(0.3), font_size=14, bold=True, color=DARK)
    add_text(slide, desc, Inches(7.1), y + Inches(0.25), Inches(5), Inches(0.3), font_size=12, color=GRAY)
    y += Inches(0.6)
add_text(slide, "DeepSeek V3 is recommended — free, fast, sufficient.", Inches(7.1), Inches(4.3), Inches(5), Inches(0.4), font_size=12, color=ACCENT, italic=True)
add_separator(slide, Inches(0.8), Inches(5.5))
add_footer(slide, "No credit card needed. Free tier is enough for this workshop.", Inches(0.8), Inches(5.8))
slide_number(slide, 8)

# ─── SLIDE 9: TEMPLATE ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_label(slide, "BOT TEMPLATE", Inches(0.8), Inches(0.5))
add_title(slide, "Template bot — 230 lines of Python.", Inches(0.8), Inches(0.95))
add_subtitle(slide, "Everything is pre-built. You just need to run it.", Inches(1.0), Inches(1.7))
# Left
add_text(slide, "trading_bot_template.py includes:", Inches(0.8), Inches(2.3), Inches(5.5), Inches(0.4), font_size=15, bold=True, color=DARK)
items = [
    "Real-time price fetching (demo mode)",
    "5 indicators: RSI, EMA, MACD, BB, ATR",
    "Signal generation with confidence scoring",
    "Auto SL/TP calculation via ATR",
    "Cooldown system to prevent spam",
    "JSON history tracking",
]
y = Inches(2.8)
for item in items:
    add_check_item(slide, item, Inches(0.8), y, font_size=14)
    y += Inches(0.38)
# Right - code snippet
add_code_block(slide, (
    "INDICATORS = {\n"
    '    "RSI": "Relative Strength Index",\n'
    '    "EMA": "Exponential Moving Average",\n'
    '    "MACD": "MACD momentum",\n'
    '    "BB": "Bollinger Bands",\n'
    '    "ATR": "Average True Range",\n'
    "}"
), Inches(6.8), Inches(2.3), Inches(5.5), Inches(1.5))
add_code_block(slide, (
    "def generate_signal(prices, history, last_signal):\n"
    '    """Combine indicators -> confidence 0-100%.\n'
    "    >= 60% -> LONG/SHORT signal.\n"
    '    < 60% -> skip (too risky)."""'
), Inches(6.8), Inches(4.0), Inches(5.5), Inches(1.1))
add_separator(slide, Inches(0.8), Inches(5.5))
add_footer(slide, "No coding required. Just download and run.", Inches(0.8), Inches(5.8))
slide_number(slide, 9)

# ─── SLIDE 10: RUN BOT ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_label(slide, "RUN BOT", Inches(0.8), Inches(0.5))
add_title(slide, "Jalanin bot dan liat sinyal pertamamu.", Inches(0.8), Inches(0.95))
add_code_block(slide, "python trading_bot_template.py", Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.45))
# Output card
card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.7), Inches(5.5), Inches(2.8))
card.fill.solid()
card.fill.fore_color.rgb = RGBColor(238, 242, 255)
card.line.fill.background()
add_text(slide, "Expected Output:", Inches(7.1), Inches(1.9), Inches(5), Inches(0.3), font_size=14, bold=True, color=ACCENT)
output_text = (
    "TRADING SIGNAL - 25 Jul 13:30 WIB\n"
    "========================================\n"
    "Price: $4,075.70 (+1.71%)\n"
    "Direction: LONG | Confidence: 95%\n"
    "Entry: $4,075.70 | SL: $3,962.50\n"
    "TP: $4,113.20 | R:R 1:1.67\n"
    "\n"
    "Reasons:\n"
    "  - RSI 66>50 bullish bias\n"
    "  - EMA9>21 bullish\n"
    "  - MACD+ bullish momentum"
)
add_text(slide, output_text, Inches(7.1), Inches(2.3), Inches(5), Inches(2.0), font_size=11, color=DARK, font_name='Consolas')
# Notes
add_text(slide, "Bot also saves history to:", Inches(0.8), Inches(3.0), Inches(5.5), Inches(0.3), font_size=14, color=GRAY)
add_text(slide, "signals.json  +  history.json", Inches(0.8), Inches(3.4), Inches(5.5), Inches(0.3), font_size=14, bold=True, color=ACCENT, font_name='Consolas')
add_text(slide, "If bot is silent — it's in cooldown or position active.", Inches(0.8), Inches(4.2), Inches(5.5), Inches(0.3), font_size=13, color=GRAY, italic=True)
add_separator(slide, Inches(0.8), Inches(5.5))
add_footer(slide, "Silent mode is a feature, not a bug. No spam.", Inches(0.8), Inches(5.8))
slide_number(slide, 10)

# ─── SLIDE 11: BOT BEHAVIOR ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_label(slide, "BOT BEHAVIOR", Inches(0.8), Inches(0.5))
add_title(slide, "Kenapa bot ga ngomong tiap detik?", Inches(0.8), Inches(0.95))
add_text(slide, "It's a FEATURE, not a bug.", Inches(0.8), Inches(1.7), Inches(8), Inches(0.4), font_size=18, bold=True, color=ACCENT)
# Cards
reasons = [
    ("Cooldown 4h", "Prevents spam. Bot waits 4 hours between signals."),
    ("Active Guard", "No new signal if a position is still open."),
    ("Silent Mode", "Output only on NEW signal or CLOSED position."),
]
for i, (title, desc) in enumerate(reasons):
    x = Inches(0.8) + i * Inches(4.1)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.3), Inches(3.8), Inches(1.8))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(238, 242, 255)
    card.line.fill.background()
    add_text(slide, title, x + Inches(0.2), Inches(2.5), Inches(3.4), Inches(0.4), font_size=16, bold=True, color=ACCENT)
    add_text(slide, desc, x + Inches(0.2), Inches(3.0), Inches(3.4), Inches(0.8), font_size=13, color=DARK)
add_text(slide, "Want to check status? Open signals.json or re-run the bot.", Inches(0.8), Inches(4.8), Inches(10), Inches(0.4), font_size=14, color=GRAY)
add_separator(slide, Inches(0.8), Inches(5.5))
add_footer(slide, "Your bot is working. It's just being polite.", Inches(0.8), Inches(5.8))
slide_number(slide, 11)

# ─── SLIDE 12: LIVE DEMO ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_label(slide, "LIVE DEMO", Inches(0.8), Inches(0.5))
add_title(slide, "Bot XAU — hidup 24/7.", Inches(0.8), Inches(0.95))
add_subtitle(slide, "This is the instructor's real bot. Running right now.", Inches(1.0), Inches(1.7))
# Left - bot card
card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.3), Inches(5.5), Inches(2.8))
card.fill.solid()
card.fill.fore_color.rgb = RGBColor(238, 242, 255)
card.line.fill.background()
add_text(slide, "XAU/USD — ACTIVE SIGNAL", Inches(1.1), Inches(2.5), Inches(5), Inches(0.4), font_size=15, bold=True, color=ACCENT)
live_data = (
    "Direction: LONG\n"
    "Entry: $4,019.00\n"
    "Current: $4,075.70 (+1.41%)\n"
    "SL: $3,962.50 | TP: $4,113.20\n"
    "R:R 1:1.67 | Duration: 36/48h\n"
    "Win Rate: 25% (1W/4L)"
)
add_text(slide, live_data, Inches(1.1), Inches(3.0), Inches(5), Inches(1.8), font_size=14, color=DARK, font_name='Consolas')
add_text(slide, "Bot tetap profit walau win rate kecil. Risk management.", Inches(0.8), Inches(5.3), Inches(5.5), Inches(0.4), font_size=13, color=ACCENT, bold=True)
# Right
add_text(slide, "Why this matters:", Inches(6.8), Inches(2.3), Inches(5.5), Inches(0.4), font_size=15, bold=True, color=DARK)
points = [
    "This is a REAL bot — not a simulation",
    "Running before you arrived today",
    "5 signals generated, 1 win 4 loss — but overall profitable",
    "Bot keeps running even after losses",
    "I never recoded. Template → Cron → Done.",
]
y = Inches(2.8)
for p in points:
    add_check_item(slide, p, Inches(6.8), y, font_size=14)
    y += Inches(0.4)
add_separator(slide, Inches(0.8), Inches(6.0))
add_footer(slide, "Your bot will look exactly like this by the end of session.", Inches(0.8), Inches(6.3))
slide_number(slide, 12)

# ─── SLIDE 13: CRON ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_label(slide, "AUTOMATION", Inches(0.8), Inches(0.5))
add_title(slide, "Cron: bot jalan sendiri, lo tinggal terima hasil.", Inches(0.8), Inches(0.95))
add_text(slide, "Running python trading_bot.py manually every day? Tiring.", Inches(0.8), Inches(1.7), Inches(10), Inches(0.4), font_size=15, color=GRAY)
add_text(slide, "Cron = bot runs itself. No server needed.", Inches(0.8), Inches(2.2), Inches(10), Inches(0.4), font_size=18, bold=True, color=BLACK)
add_code_block(slide, (
    "hermes cron create 'every 4h' \\\n"
    '    --name "Sinyal XAU Auto" \\\n'
    '    --prompt "Jalankan python trading_bot_template.py 2>&1. Kirim output apa adanya."'
), Inches(0.8), Inches(2.8), Inches(11), Inches(1.2))
add_text(slide, "Check running crons:", Inches(0.8), Inches(4.3), Inches(5), Inches(0.4), font_size=14, color=GRAY)
add_code_block(slide, "hermes cron list", Inches(0.8), Inches(4.6), Inches(4), Inches(0.4))
add_text(slide, "You now have a 24/7 trading signal bot. Congratulations!", Inches(0.8), Inches(5.3), Inches(10), Inches(0.4), font_size=16, bold=True, color=ACCENT)
add_separator(slide, Inches(0.8), Inches(5.8))
add_footer(slide, "Cron handles everything. You just read the results.", Inches(0.8), Inches(6.1))
slide_number(slide, 13)

# ─── SLIDE 14: TELEGRAM ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_label(slide, "NOTIFICATIONS", Inches(0.8), Inches(0.5))
add_title(slide, "Bonus: Kirim sinyal ke Telegram.", Inches(0.8), Inches(0.95))
# Steps
steps = [
    ("Create bot via @BotFather", "Get your bot token"),
    ("Start your bot on phone", "Search bot name, click start"),
    ("Setup Hermes Gateway", "hermes gateway setup → select Telegram"),
    ("Update cron delivery", "hermes cron update <id> --deliver telegram"),
]
for i, (title, desc) in enumerate(steps):
    add_step_item(slide, i+1, title, desc, Inches(0.8), Inches(1.7) + i * Inches(0.6))
# Right - result card
card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.3), Inches(1.7), Inches(5), Inches(3.5))
card.fill.solid()
card.fill.fore_color.rgb = RGBColor(238, 242, 255)
card.line.fill.background()
add_text(slide, "Result:", Inches(7.6), Inches(1.9), Inches(4.5), Inches(0.4), font_size=16, bold=True, color=ACCENT)
msg = "[TRADING SIGNAL - 25 Jul 14:00]\n"
msg += "Price: $4,075.70 (+1.71%)\n"
msg += "Direction: LONG | Confidence: 95%\n"
msg += "Entry: $4,075.70 | R:R 1:1.67\n"
msg += "\n"
msg += "Reasons:\n"
msg += " - RSI bullish, EMA9>21\n"
msg += " - MACD+ momentum"
add_text(slide, msg, Inches(7.6), Inches(2.4), Inches(4.5), Inches(1.8), font_size=12, color=DARK, font_name='Consolas')
add_text(slide, "Every 4 hours, signal arrives on your phone.", Inches(7.6), Inches(4.5), Inches(4.5), Inches(0.4), font_size=14, bold=True, color=ACCENT)
add_text(slide, "No laptop needed after setup.", Inches(7.6), Inches(4.9), Inches(4.5), Inches(0.4), font_size=13, color=GRAY)
add_separator(slide, Inches(0.8), Inches(5.8))
add_footer(slide, "Telegram is optional but highly recommended.", Inches(0.8), Inches(6.1))
slide_number(slide, 14)

# ─── SLIDE 15: HANDS-ON ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_label(slide, "HANDS-ON", Inches(0.8), Inches(0.5))
add_title(slide, "Practice time — 15 minutes.", Inches(0.8), Inches(0.95))
add_subtitle(slide, "Follow along. Help your neighbor if they're stuck.", Inches(1.0), Inches(1.7))
steps = [
    ("Install Hermes", "curl -fsSL https://hermes-agent.nousresearch.com/install.sh | bash"),
    ("Setup provider", "hermes model → OpenRouter → paste key → DeepSeek V3"),
    ("Download template", "QR code / link provided"),
    ("Run bot", "python trading_bot_template.py"),
    ("Setup cron", "hermes cron create 'every 6h' --name 'Sinyal Saya' --prompt ..."),
    ("Telegram (optional)", "hermes gateway setup"),
]
for i, (title, cmd) in enumerate(steps):
    col = 0 if i < 3 else 1
    row = i if i < 3 else i - 3
    x = Inches(0.8) if col == 0 else Inches(6.8)
    y = Inches(2.2) + row * Inches(0.8)
    add_step_item(slide, i+1, title, cmd, x, y, width=Inches(5.5))
add_text(slide, "Done! Your trading bot is now live.", Inches(0.8), Inches(5.8), Inches(10), Inches(0.5), font_size=18, bold=True, color=ACCENT)
slide_number(slide, 15)

# ─── SLIDE 16: TIPS ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_label(slide, "BEST PRACTICES", Inches(0.8), Inches(0.5))
add_title(slide, "Tips biar trading bot lo sukses.", Inches(0.8), Inches(0.95))
# Left - DO
add_text(slide, "DO:", Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.4), font_size=16, bold=True, color=GREEN)
do_items = [
    "Start with 1 asset (XAU or BTC)",
    "Always set Stop Loss",
    "Backtest 1 week before real trade",
    "Cross-check AI signals with chart",
    "Be consistent — 3 good signals > 20 random",
]
y = Inches(2.2)
for item in do_items:
    add_check_item(slide, item, Inches(0.8), y, font_size=14, color=DARK)
    y += Inches(0.38)
# Right - DON'T
add_text(slide, "DON'T:", Inches(6.8), Inches(1.7), Inches(5.5), Inches(0.4), font_size=16, bold=True, color=RED)
dont_items = [
    "Trade real money immediately",
    "Over-trade (cron every hour = spam)",
    "Trust AI 100% — you're still the boss",
]
y = Inches(2.2)
for item in dont_items:
    add_check_item(slide, item, Inches(6.8), y, font_size=14, color=RED)
    y += Inches(0.4)
add_text(slide, "Risk management > Win rate. Always.", Inches(0.8), Inches(4.8), Inches(10), Inches(0.4), font_size=16, bold=True, color=ACCENT)
add_separator(slide, Inches(0.8), Inches(5.5))
add_footer(slide, "Start small, learn, then scale.", Inches(0.8), Inches(5.8))
slide_number(slide, 16)

# ─── SLIDE 17: RESOURCES ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_label(slide, "RESOURCES", Inches(0.8), Inches(0.5))
add_title(slide, "All materials — scan QR or use links.", Inches(0.8), Inches(0.95))
resources = [
    ("Hermes Agent", "hermes-agent.nousresearch.com"),
    ("OpenRouter (Free API Key)", "openrouter.ai/keys"),
    ("Python Download", "python.org/downloads"),
    ("Bot Template File", "[QR Code Provided]"),
    ("Cheatsheet Prompts", "[QR Code Provided]"),
    ("PDF GuideBook", "[QR Code Provided]"),
]
for i, (title, url) in enumerate(resources):
    col = 0 if i < 3 else 1
    row = i if i < 3 else i - 3
    x = Inches(0.8) if col == 0 else Inches(6.8)
    y = Inches(1.7) + row * Inches(0.8)
    add_text(slide, title, x, y, Inches(5.5), Inches(0.35), font_size=15, bold=True, color=DARK)
    add_text(slide, url, x, y + Inches(0.35), Inches(5.5), Inches(0.35), font_size=13, color=ACCENT)
add_text(slide, "All QR codes are in the PDF GuideBook — available after session.", Inches(0.8), Inches(5.2), Inches(10), Inches(0.4), font_size=13, color=GRAY, italic=True)
add_separator(slide, Inches(0.8), Inches(5.6))
add_footer(slide, "You get everything to continue learning at home.", Inches(0.8), Inches(5.9))
slide_number(slide, 17)

# ─── SLIDE 18: Q&A ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_text(slide, "Q & A", Inches(0), Inches(2.0), Inches(13.333), Inches(1.2), font_size=56, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_separator(slide, Inches(4.5), Inches(3.4), Inches(4.333))
add_text(slide, "Thank you!", Inches(0), Inches(3.6), Inches(13.333), Inches(0.7), font_size=32, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
add_text(slide, "You now have your own AI trading signal generator.", Inches(0), Inches(4.6), Inches(13.333), Inches(0.5), font_size=18, color=DARK, align=PP_ALIGN.CENTER)
add_text(slide, "No coding experience needed. Just prompt, download, cron, done.", Inches(0), Inches(5.2), Inches(13.333), Inches(0.5), font_size=16, color=ACCENT, align=PP_ALIGN.CENTER, italic=True)
add_text(slide, "@Putra", Inches(0), Inches(6.2), Inches(13.333), Inches(0.4), font_size=14, color=GRAY, align=PP_ALIGN.CENTER)
slide_number(slide, 18)

# ─── SAVE ───
out_path = os.path.expanduser('~/Desktop/workshop-pintu/FROM-ZERO-TO-TRADING-BOTS-v2.pptx')
prs.save(out_path)
print(f"PPTX saved: {out_path}")
print(f"Size: {os.path.getsize(out_path)} bytes")
print(f"Slides: {len(prs.slides)}")
